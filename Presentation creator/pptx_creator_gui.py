import sys
import json
import os
from PyQt5.QtWidgets import (QApplication, QMainWindow, QVBoxLayout, QHBoxLayout, 
                             QWidget, QPushButton, QLabel, QLineEdit, QTextEdit, 
                             QSpinBox, QListWidget, QListWidgetItem, QMessageBox,
                             QFileDialog, QTabWidget, QFormLayout, QGroupBox)
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QFont, QIcon
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt

class SlideEditor(QWidget):
    def __init__(self):
        super().__init__()
        self.init_ui()
        
    def init_ui(self):
        layout = QVBoxLayout()
        
        # Title input
        title_layout = QHBoxLayout()
        title_layout.addWidget(QLabel("Title:"))
        self.title_edit = QLineEdit()
        title_layout.addWidget(self.title_edit)
        layout.addLayout(title_layout)
        
        # Content input
        layout.addWidget(QLabel("Content:"))
        self.content_edit = QTextEdit()
        self.content_edit.setMaximumHeight(150)
        layout.addWidget(self.content_edit)
        
        self.setLayout(layout)
    
    def get_slide_data(self):
        return {
            "title": self.title_edit.text(),
            "content": self.content_edit.toPlainText()
        }
    
    def set_slide_data(self, title, content):
        self.title_edit.setText(title)
        self.content_edit.setPlainText(content)
    
    def clear(self):
        self.title_edit.clear()
        self.content_edit.clear()

class PowerPointGUI(QMainWindow):
    def __init__(self):
        super().__init__()
        self.slides_data = []
        self.current_slide_index = -1
        self.init_ui()
        
    def init_ui(self):
        self.setWindowTitle("PowerPoint Presentation Creator")
        self.setGeometry(100, 100, 800, 600)
        
        # Create central widget and main layout
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QHBoxLayout(central_widget)
        
        # Left panel - Slide list and controls
        left_panel = self.create_left_panel()
        main_layout.addWidget(left_panel, 1)
        
        # Right panel - Slide editor and settings
        right_panel = self.create_right_panel()
        main_layout.addWidget(right_panel, 2)
        
    def create_left_panel(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        
        # Slide list
        layout.addWidget(QLabel("Slides:"))
        self.slides_list = QListWidget()
        self.slides_list.itemClicked.connect(self.on_slide_selected)
        layout.addWidget(self.slides_list)
        
        # Slide control buttons
        buttons_layout = QVBoxLayout()
        
        self.add_slide_btn = QPushButton("Add Slide")
        self.add_slide_btn.clicked.connect(self.add_slide)
        buttons_layout.addWidget(self.add_slide_btn)
        
        self.remove_slide_btn = QPushButton("Remove Slide")
        self.remove_slide_btn.clicked.connect(self.remove_slide)
        buttons_layout.addWidget(self.remove_slide_btn)
        
        self.move_up_btn = QPushButton("Move Up")
        self.move_up_btn.clicked.connect(self.move_slide_up)
        buttons_layout.addWidget(self.move_up_btn)
        
        self.move_down_btn = QPushButton("Move Down")
        self.move_down_btn.clicked.connect(self.move_slide_down)
        buttons_layout.addWidget(self.move_down_btn)
        
        layout.addLayout(buttons_layout)
        
        # JSON operations
        json_group = QGroupBox("JSON Operations")
        json_layout = QVBoxLayout(json_group)
        
        self.load_json_btn = QPushButton("Load from JSON")
        self.load_json_btn.clicked.connect(self.load_from_json)
        json_layout.addWidget(self.load_json_btn)
        
        self.save_json_btn = QPushButton("Save to JSON")
        self.save_json_btn.clicked.connect(self.save_to_json)
        json_layout.addWidget(self.save_json_btn)
        
        layout.addWidget(json_group)
        
        return panel
    
    def create_right_panel(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        
        # Tab widget for editor and settings
        tab_widget = QTabWidget()
        
        # Slide Editor Tab
        editor_tab = QWidget()
        editor_layout = QVBoxLayout(editor_tab)
        
        editor_layout.addWidget(QLabel("Edit Current Slide:"))
        self.slide_editor = SlideEditor()
        editor_layout.addWidget(self.slide_editor)
        
        # Update slide button
        self.update_slide_btn = QPushButton("Update Slide")
        self.update_slide_btn.clicked.connect(self.update_current_slide)
        editor_layout.addWidget(self.update_slide_btn)
        
        tab_widget.addTab(editor_tab, "Slide Editor")
        
        # Settings Tab
        settings_tab = self.create_settings_tab()
        tab_widget.addTab(settings_tab, "Font Settings")
        
        layout.addWidget(tab_widget)
        
        # Presentation operations
        ppt_group = QGroupBox("Presentation Operations")
        ppt_layout = QVBoxLayout(ppt_group)
        
        # Filename input
        filename_layout = QHBoxLayout()
        filename_layout.addWidget(QLabel("Filename:"))
        self.filename_edit = QLineEdit("presentation")
        filename_layout.addWidget(self.filename_edit)
        ppt_layout.addLayout(filename_layout)
        
        # Create and modify buttons
        self.create_ppt_btn = QPushButton("Create New Presentation")
        self.create_ppt_btn.clicked.connect(self.create_presentation)
        ppt_layout.addWidget(self.create_ppt_btn)
        
        self.modify_ppt_btn = QPushButton("Modify Existing Presentation")
        self.modify_ppt_btn.clicked.connect(self.modify_presentation)
        ppt_layout.addWidget(self.modify_ppt_btn)
        
        layout.addWidget(ppt_group)
        
        return panel
    
    def create_settings_tab(self):
        settings_tab = QWidget()
        layout = QFormLayout(settings_tab)
        
        # Font size settings
        self.title_font_size = QSpinBox()
        self.title_font_size.setRange(8, 72)
        self.title_font_size.setValue(36)
        self.title_font_size.setSuffix(" pt")
        layout.addRow("Title Font Size:", self.title_font_size)
        
        self.content_font_size = QSpinBox()
        self.content_font_size.setRange(8, 72)
        self.content_font_size.setValue(32)
        self.content_font_size.setSuffix(" pt")
        layout.addRow("Content Font Size:", self.content_font_size)
        
        return settings_tab
    
    def add_slide(self):
        slide_data = {"title": f"Slide {len(self.slides_data) + 1}", "content": "Enter content here"}
        self.slides_data.append(slide_data)
        self.update_slides_list()
        self.slides_list.setCurrentRow(len(self.slides_data) - 1)
        self.load_slide_to_editor(len(self.slides_data) - 1)
    
    def remove_slide(self):
        current_row = self.slides_list.currentRow()
        if current_row >= 0 and current_row < len(self.slides_data):
            self.slides_data.pop(current_row)
            self.update_slides_list()
            self.slide_editor.clear()
            self.current_slide_index = -1
    
    def move_slide_up(self):
        current_row = self.slides_list.currentRow()
        if current_row > 0:
            self.slides_data[current_row], self.slides_data[current_row - 1] = \
                self.slides_data[current_row - 1], self.slides_data[current_row]
            self.update_slides_list()
            self.slides_list.setCurrentRow(current_row - 1)
    
    def move_slide_down(self):
        current_row = self.slides_list.currentRow()
        if current_row >= 0 and current_row < len(self.slides_data) - 1:
            self.slides_data[current_row], self.slides_data[current_row + 1] = \
                self.slides_data[current_row + 1], self.slides_data[current_row]
            self.update_slides_list()
            self.slides_list.setCurrentRow(current_row + 1)
    
    def update_slides_list(self):
        self.slides_list.clear()
        for i, slide in enumerate(self.slides_data):
            title = slide.get("title", f"Slide {i + 1}")
            item = QListWidgetItem(f"{i + 1}. {title}")
            self.slides_list.addItem(item)
    
    def on_slide_selected(self, item):
        row = self.slides_list.row(item)
        self.load_slide_to_editor(row)
    
    def load_slide_to_editor(self, index):
        if 0 <= index < len(self.slides_data):
            self.current_slide_index = index
            slide = self.slides_data[index]
            self.slide_editor.set_slide_data(slide["title"], slide["content"])
    
    def update_current_slide(self):
        if self.current_slide_index >= 0:
            slide_data = self.slide_editor.get_slide_data()
            self.slides_data[self.current_slide_index] = slide_data
            self.update_slides_list()
            self.slides_list.setCurrentRow(self.current_slide_index)
    
    def load_from_json(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Load JSON File", "", "JSON Files (*.json)")
        if file_path:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    self.slides_data = json.load(f)
                self.update_slides_list()
                if self.slides_data:
                    self.slides_list.setCurrentRow(0)
                    self.load_slide_to_editor(0)
                QMessageBox.information(self, "Success", f"Loaded {len(self.slides_data)} slides from JSON")
            except Exception as e:
                QMessageBox.critical(self, "Error", f"Failed to load JSON: {str(e)}")
    
    def save_to_json(self):
        if not self.slides_data:
            QMessageBox.warning(self, "Warning", "No slides to save")
            return
        
        file_path, _ = QFileDialog.getSaveFileName(self, "Save JSON File", "slides.json", "JSON Files (*.json)")
        if file_path:
            try:
                with open(file_path, 'w', encoding='utf-8') as f:
                    json.dump(self.slides_data, f, indent=2, ensure_ascii=False)
                QMessageBox.information(self, "Success", f"Saved {len(self.slides_data)} slides to JSON")
            except Exception as e:
                QMessageBox.critical(self, "Error", f"Failed to save JSON: {str(e)}")
    
    def change_font_and_size(self, presentation, content_font_size, title_font_size):
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if shape == slide.shapes.title:
                                run.font.name = 'Times New Roman'
                                run.font.size = title_font_size
                                paragraph.alignment = PP_ALIGN.CENTER
                                run.font.color.rgb = RGBColor(0, 0, 0)
                            else:
                                run.font.name = 'Times New Roman'
                                run.font.size = content_font_size
                                paragraph.alignment = PP_ALIGN.LEFT
                                run.font.color.rgb = RGBColor(0, 0, 0)
    
    def create_presentation_from_data(self, slides_content):
        presentation = Presentation()
        for slide in slides_content:
            slide_layout = presentation.slide_layouts[1]
            slide_obj = presentation.slides.add_slide(slide_layout)
            title_placeholder = slide_obj.shapes.title
            body_placeholder = slide_obj.placeholders[1]
            title_placeholder.text = slide['title']
            tf = body_placeholder.text_frame
            tf.text = slide['content']
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Times New Roman'
                    run.font.size = Pt(self.content_font_size.value())
        return presentation
    
    def create_presentation(self):
        if not self.slides_data:
            QMessageBox.warning(self, "Warning", "No slides to create presentation")
            return
        
        # Update current slide if being edited
        if self.current_slide_index >= 0:
            self.update_current_slide()
        
        filename = self.filename_edit.text().strip()
        if not filename:
            filename = "presentation"
        
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        try:
            presentation = self.create_presentation_from_data(self.slides_data)
            self.change_font_and_size(
                presentation, 
                Pt(self.content_font_size.value()), 
                Pt(self.title_font_size.value())
            )
            presentation.save(filename)
            QMessageBox.information(self, "Success", f"Presentation saved as {filename}")
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to create presentation: {str(e)}")
    
    def modify_presentation(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Select PowerPoint File", "", "PowerPoint Files (*.pptx)")
        if file_path:
            try:
                presentation = Presentation(file_path)
                self.change_font_and_size(
                    presentation, 
                    Pt(self.content_font_size.value()), 
                    Pt(self.title_font_size.value())
                )
                
                # Create modified filename
                base_name = os.path.splitext(file_path)[0]
                modified_filename = f"{base_name}_modified.pptx"
                presentation.save(modified_filename)
                
                QMessageBox.information(self, "Success", 
                    f"Font and size changed to Times New Roman and saved as {os.path.basename(modified_filename)}")
            except Exception as e:
                QMessageBox.critical(self, "Error", f"Failed to modify presentation: {str(e)}")

def main():
    app = QApplication(sys.argv)
    app.setStyle('Fusion')  # Modern look
    
    window = PowerPointGUI()
    """
    # Sample JSON structure for slides data (should be saved as .json file):
    window.slides_data = [
        {"title": "Welcome", "content": "This is the first slide\nAdd your content here"},
        {"title": "Overview", "content": "• Point 1\n• Point 2\n• Point 3"},
        {"title": "Conclusion", "content": "Thank you for your attention!"}
    ]
    """
    
    window.show()
    sys.exit(app.exec_())

if __name__ == "__main__":
    main()